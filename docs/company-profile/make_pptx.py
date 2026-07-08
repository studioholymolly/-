# -*- coding: utf-8 -*-
"""스튜디오 홀리몰리 회사소개서 PPTX v2 (구글 슬라이드 편집용).

레퍼런스: Pentagram / COLLINS / Koto 계열 credentials deck 문법 —
거대한 디스플레이 타이포 vs 소형 메타 텍스트의 스케일 대비,
페이지당 하나의 메시지, 블랙/페이퍼 교차 리듬, 풀블리드 이미지.
폰트는 구글 폰트인 Jost(영문 디스플레이) + Noto Sans KR(국문)만 사용.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

PAPER = RGBColor(0xFA, 0xFA, 0xF8)
INK = RGBColor(0x0A, 0x0A, 0x0A)
GRAY = RGBColor(0x6F, 0x6F, 0x6F)
GRAY3 = RGBColor(0xA3, 0xA3, 0xA3)
LINE = RGBColor(0xE5, 0xE5, 0xE3)
DARK = RGBColor(0x11, 0x11, 0x11)
DGRAY = RGBColor(0x9A, 0x9A, 0x9A)   # dark-slide body
DNUM = RGBColor(0x4A, 0x4A, 0x4A)    # dark-slide giant numbers
DLINE = RGBColor(0x2E, 0x2E, 0x2E)   # dark-slide hairline
DLABEL = RGBColor(0x8A, 0x8A, 0x8A)  # dark-slide label
PH_FILL = RGBColor(0xEF, 0xEF, 0xED)  # image placeholder

DISPLAY = 'Jost'
KR = 'Noto Sans KR'

W, H = Inches(13.333), Inches(7.5)
MX = Inches(0.66)
CW = W - MX * 2

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def set_run(run, size, color, bold=False, latin=DISPLAY, ea=KR, spacing=None):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = latin
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:ea', 'a:cs'):
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
    latin_el = rPr.find(qn('a:latin'))
    ea_el = rPr.makeelement(qn('a:ea'), {'typeface': ea})
    cs_el = rPr.makeelement(qn('a:cs'), {'typeface': ea})
    latin_el.addnext(cs_el)
    latin_el.addnext(ea_el)
    if spacing is not None:
        rPr.set('spc', str(spacing))


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
            line_spacing=None, space_after=None, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (text, kw) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = text
        set_run(run, **kw)
    return box


def rect(slide, x, y, w, h, fill, outline=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if outline is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = outline
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    style = sh._element.find(qn('p:style'))
    if style is not None:
        sh._element.remove(style)
    return sh


def hairline(slide, x, y, w, color=LINE, weight=1.0):
    return rect(slide, x, y, w, Emu(int(Pt(weight))), color)


def new_slide(bg=PAPER):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, W, H, bg)
    return slide


def meta(slide, label_text, num, dark=False):
    """모든 페이지 공통 메타 바: 라벨 — 헤어라인 — 페이지 번호."""
    lc = DLABEL if dark else GRAY
    line_c = DLINE if dark else LINE
    textbox(slide, MX, Inches(0.34), Inches(9), Inches(0.3),
            [(label_text.upper(), dict(size=9, color=lc, spacing=220))])
    textbox(slide, W - MX - Inches(1.2), Inches(0.34), Inches(1.2), Inches(0.3),
            [(f'{num:02d}', dict(size=9, color=lc, spacing=220))], align=PP_ALIGN.RIGHT)
    hairline(slide, MX, Inches(0.72), CW, line_c)


def img_placeholder(slide, x, y, w, h, note, border=True):
    rect(slide, x, y, w, h, PH_FILL, outline=LINE if border else None)
    textbox(slide, x, y, w, h,
            [('IMAGE', dict(size=10, color=GRAY3, spacing=280)),
             (note, dict(size=9.5, color=GRAY3, latin=KR))],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=4)


# ================================================== 01 — 커버
s = new_slide()
textbox(s, MX, Inches(0.34), Inches(10), Inches(0.3),
        [('STUDIO HOLYMOLLY — VISUAL DIRECTION STUDIO', dict(size=9, color=GRAY, spacing=220))])
textbox(s, W - MX - Inches(4), Inches(0.34), Inches(4), Inches(0.3),
        [('COMPANY PROFILE 2026', dict(size=9, color=GRAY, spacing=220))], align=PP_ALIGN.RIGHT)
hairline(s, MX, Inches(0.72), CW)
textbox(s, MX - Inches(0.04), Inches(1.75), Inches(13.4), Inches(4.2),
        [('STUDIO.', dict(size=130, color=INK, spacing=150)),
         ('HOLYMOLLY', dict(size=130, color=INK, spacing=150))],
        line_spacing=0.98)
hairline(s, MX, Inches(6.45), CW)
textbox(s, MX, Inches(6.68), Inches(7), Inches(0.5),
        [('브랜드가 보여지는 모든 장면을 만듭니다.', dict(size=12.5, color=INK, bold=True, latin=KR))])
textbox(s, W - MX - Inches(6), Inches(6.7), Inches(6), Inches(0.5),
        [('SEOUL, YEOKSAM — STUDIOHOLYMOLLY.MYPORTFOLIO.COM', dict(size=8.5, color=GRAY3, spacing=150))],
        align=PP_ALIGN.RIGHT)

# ================================================== 02 — 매니페스토 (블랙)
s = new_slide(DARK)
textbox(s, MX, Inches(0.34), Inches(9), Inches(0.3),
        [('PERSPECTIVE', dict(size=9, color=DLABEL, spacing=220))])
textbox(s, MX, Inches(2.2), Inches(12.2), Inches(3.6),
        [('컷 한 장을 찍는 일보다,', dict(size=44, color=PAPER, bold=True, latin=KR)),
         ('그 컷이 브랜드 안에서 어떻게', dict(size=44, color=PAPER, bold=True, latin=KR)),
         ('보일지를 먼저 생각합니다.', dict(size=44, color=PAPER, bold=True, latin=KR))],
        line_spacing=1.32)
textbox(s, W - MX - Inches(1.2), H - Inches(0.6), Inches(1.2), Inches(0.3),
        [('02', dict(size=9, color=DLABEL, spacing=220))], align=PP_ALIGN.RIGHT)

# ================================================== 03 — Who We Are + 팩트 스트립
s = new_slide()
meta(s, 'Who We Are', 3)
textbox(s, MX, Inches(1.45), Inches(7.2), Inches(2.2),
        [('서울 역삼,', dict(size=40, color=INK, bold=True, latin=KR)),
         ('비주얼 디렉션 스튜디오.', dict(size=40, color=INK, bold=True, latin=KR))],
        line_spacing=1.28)
textbox(s, Inches(8.35), Inches(1.62), W - MX - Inches(8.35), Inches(3.3),
        [('스튜디오 홀리몰리는 사진과 영상, 그리고 BX 디자인까지 — 하나의 무드로 이어지는 장면을 설계합니다.',
          dict(size=11, color=GRAY, latin=KR)),
         ('뷰티와 F&B, 화면 속 질감이 곧 구매가 되는 카테고리에서 출발해 제품 · 라이프스타일 · 패션 · 인물로 영역을 넓혀왔습니다. 그렇게 찍은 컷이 상세페이지가 되고, 캠페인이 되고, 브랜드의 얼굴이 됩니다.',
          dict(size=11, color=GRAY, latin=KR))],
        line_spacing=1.6, space_after=12)
facts = [
    ('FIELD', 'Photo · Film · Branding', 13, DISPLAY),
    ('CATEGORY', 'Beauty · F&B · Lifestyle', 13, DISPLAY),
    ('LOCATION', '역삼동 783-2, B1', 13, KR),
    ('CONTACT', 'studio.holymolly@gmail.com', 10, DISPLAY),
]
fy = Inches(5.35)
hairline(s, MX, fy, CW, INK, 1.3)
col_w = CW / 4
for i, (lab, val, vs, lat) in enumerate(facts):
    x = MX + col_w * i
    if i > 0:
        rect(s, x - Emu(int(Pt(0.45))), fy, Emu(int(Pt(0.9))), Inches(1.5), LINE)
    pad = Inches(0.3) if i else 0
    textbox(s, x + pad, fy + Inches(0.3), col_w - Inches(0.55), Inches(1.1),
            [(lab, dict(size=8.5, color=GRAY, spacing=220)),
             (val, dict(size=vs, color=INK, bold=True, latin=lat))],
            line_spacing=1.75)
hairline(s, MX, fy + Inches(1.5), CW)

# ================================================== 04 — Services 인덱스 (한 페이지)
services = [
    ('01', '뷰티 & 제품', '제형과 컬러, 형태와 소재를 정교하게. 커머스 컷부터 캠페인까지, 제품이 팔리게 하는 컷.'),
    ('02', 'F&B & 라이프스타일', '시즐과 온도, 쓰는 장면과 먹는 순간. 제품이 놓일 실제 삶의 장면을 만듭니다.'),
    ('03', '패션 & 인물', '시즌 룩북과 커머스 화보, 브랜드를 이끄는 사람들. 브랜드의 결을 인물로 옮깁니다.'),
    ('04', '영상', '브랜드 필름, 커머스 영상, 소셜 콘텐츠. 같은 무드를 무빙으로 확장합니다.'),
    ('05', 'BX 디자인', '비주얼 아이덴티티부터 콘텐츠 시스템까지. 장면이 브랜드가 되도록.'),
]
s = new_slide()
meta(s, 'What We Do', 4)
y = Inches(1.25)
row_h = Inches(1.06)
for num, title, body in services:
    if y > Inches(1.3):
        hairline(s, MX, y, CW)
    textbox(s, MX, y + Inches(0.28), Inches(1.3), Inches(0.6),
            [(num, dict(size=24, color=GRAY3, spacing=100))])
    textbox(s, MX + Inches(1.55), y + Inches(0.24), Inches(5.6), Inches(0.7),
            [(title, dict(size=23, color=INK, bold=True, latin=KR))])
    textbox(s, Inches(7.55), y + Inches(0.34), W - MX - Inches(7.55), Inches(0.7),
            [(body, dict(size=9.5, color=GRAY, latin=KR))], line_spacing=1.45)
    y += row_h
hairline(s, MX, y, CW, INK, 1.3)
textbox(s, MX, y + Inches(0.22), CW, Inches(0.4),
        [('사진 · 영상 · 디자인을 한 팀이 진행하면, 채널 전체가 하나의 무드로 이어집니다.',
          dict(size=10.5, color=GRAY, latin=KR))])

# ================================================== 05 — WORKS 디바이더 (블랙)
s = new_slide(DARK)
textbox(s, MX, Inches(0.34), Inches(9), Inches(0.3),
        [('SELECTED WORKS', dict(size=9, color=DLABEL, spacing=220))])
textbox(s, MX - Inches(0.04), Inches(1.7), Inches(13.5), Inches(2.8),
        [('WORKS →', dict(size=160, color=PAPER, spacing=150))])
textbox(s, MX, Inches(4.55), Inches(10), Inches(0.6),
        [('말보다 컷이 빠릅니다. 작업물은 여기서 보세요.', dict(size=13, color=DGRAY, latin=KR))])
ly = Inches(5.55)
hairline(s, MX, ly, CW, DLINE)
links = [
    ('PORTFOLIO', 'studioholymolly.myportfolio.com'),
    ('INSTAGRAM', '@studio_holymolly'),
    ('THREADS', '@studio_holymolly'),
]
col_w = CW / 3
for i, (lab, val) in enumerate(links):
    x = MX + col_w * i
    textbox(s, x, ly + Inches(0.3), col_w - Inches(0.4), Inches(1.0),
            [(lab, dict(size=8.5, color=DLABEL, spacing=220)),
             (val, dict(size=13, color=PAPER, bold=True))],
            line_spacing=1.8)
textbox(s, W - MX - Inches(1.2), H - Inches(0.6), Inches(1.2), Inches(0.3),
        [('05', dict(size=9, color=DLABEL, spacing=220))], align=PP_ALIGN.RIGHT)

# ================================================== 06–09 — 케이스 (풀블리드)
cases = [
    ('CASE 01', '뷰티 대표 컷으로 교체 (풀블리드)'),
    ('CASE 02', '제품 / 커머스 대표 컷으로 교체 (풀블리드)'),
    ('CASE 03', 'F&B 대표 컷으로 교체 (풀블리드)'),
    ('CASE 04', '라이프스타일 · 인물 대표 컷으로 교체 (풀블리드)'),
]
bar_h = Inches(1.0)
for i, (case_no, img_note) in enumerate(cases):
    s = new_slide()
    img_placeholder(s, 0, 0, W, H - bar_h, img_note, border=False)
    rect(s, 0, H - bar_h, W, bar_h, PAPER)
    hairline(s, 0, H - bar_h, W, INK, 1.3)
    textbox(s, MX, H - bar_h + Inches(0.2), Inches(5.6), Inches(0.7),
            [('[클라이언트/브랜드명] — [프로젝트명]', dict(size=13, color=INK, bold=True, latin=KR)),
             ('SCOPE — 기획 · 촬영 · 리터칭 (해당 항목만)', dict(size=8.5, color=GRAY, latin=KR, spacing=80))],
            line_spacing=1.5)
    textbox(s, Inches(6.6), H - bar_h + Inches(0.24), Inches(4.6), Inches(0.6),
            [('[한 줄 요약 — 무엇을 해결했는지]', dict(size=9.5, color=GRAY, latin=KR))],
            line_spacing=1.45)
    textbox(s, W - MX - Inches(1.8), H - bar_h + Inches(0.24), Inches(1.8), Inches(0.3),
            [(f'{case_no} — {6 + i:02d}', dict(size=8.5, color=GRAY3, spacing=150))],
            align=PP_ALIGN.RIGHT)

# ================================================== 10 — Why Holymolly
strengths = [
    ('01', '디렉션이 있는 촬영', '찍기 전에 설계합니다. 레퍼런스부터 무드보드, 세팅 계획까지 — 촬영은 기획의 마지막 단계입니다.'),
    ('02', '질감에 강한 스튜디오', '뷰티의 제형, F&B의 시즐, 패브릭의 결. 클로즈업에서 무너지지 않는 디테일.'),
    ('03', '현장에서 한 컷 더', '계획에 없던 세팅을 그 자리에서 추가하는 팀. 같은 예산에서 얻는 컷이 늘어납니다.'),
    ('04', '원본부터 다른 퀄리티', '톤 정리만 거친 원본을 그대로 보여드릴 수 있는 촬영. 보정은 마지막 손질입니다.'),
]
s = new_slide()
meta(s, 'Why Holymolly', 10)
textbox(s, MX, Inches(1.35), Inches(11), Inches(0.9),
        [('홀리몰리와 찍으면 다른 것', dict(size=34, color=INK, bold=True, latin=KR))])
gy = Inches(2.75)
cell_w, cell_h = CW / 2, Inches(2.0)
hairline(s, MX, gy, CW, INK, 1.3)
hairline(s, MX, gy + cell_h, CW)
hairline(s, MX, gy + cell_h * 2, CW)
rect(s, MX + cell_w - Emu(int(Pt(0.45))), gy, Emu(int(Pt(0.9))), cell_h * 2, LINE)
for i, (num, title, body) in enumerate(strengths):
    r, c = divmod(i, 2)
    x = MX + cell_w * c + (Inches(0.35) if c else 0)
    yy = gy + cell_h * r
    textbox(s, x, yy + Inches(0.24), Inches(1.6), Inches(1.4),
            [(num, dict(size=38, color=GRAY3, spacing=50))])
    textbox(s, x + Inches(1.35), yy + Inches(0.32), cell_w - Inches(2.05), cell_h - Inches(0.5),
            [(title, dict(size=15.5, color=INK, bold=True, latin=KR)),
             (body, dict(size=10, color=GRAY, latin=KR))],
            line_spacing=1.45, space_after=6)

# ================================================== 11 — Process (블랙)
process = [
    ('01', '상담 · 기획', '24시간 안에 답장. 목적과 무드를 듣고 레퍼런스 · 세팅 · 견적을 제안합니다.'),
    ('02', '촬영', '역삼 스튜디오 또는 로케이션. 현장에서 컷을 함께 확인하며 진행합니다.'),
    ('03', '셀렉 · 리터칭', '정리된 컷에서 고르시면, 피드백을 반영해 정성껏 보정합니다.'),
    ('04', '전달', '고해상도 완성본을 일정에 맞춰. 채널별 포맷 정리까지 함께.'),
]
s = new_slide(DARK)
textbox(s, MX, Inches(0.34), Inches(9), Inches(0.3),
        [('HOW WE WORK', dict(size=9, color=DLABEL, spacing=220))])
hairline(s, MX, Inches(0.72), CW, DLINE)
textbox(s, MX, Inches(1.35), Inches(11), Inches(0.9),
        [('상담부터 전달까지, 네 걸음', dict(size=34, color=PAPER, bold=True, latin=KR))])
py = Inches(3.0)
col_w = CW / 4
for i, (num, title, body) in enumerate(process):
    x = MX + col_w * i
    if i > 0:
        rect(s, x - Emu(int(Pt(0.45))), py, Emu(int(Pt(0.9))), Inches(3.3), DLINE)
    pad = Inches(0.32) if i else 0
    textbox(s, x + pad, py + Inches(0.15), col_w - Inches(0.6), Inches(1.1),
            [(num, dict(size=44, color=DNUM, spacing=50))])
    textbox(s, x + pad, py + Inches(1.35), col_w - Inches(0.6), Inches(1.9),
            [(title, dict(size=15, color=PAPER, bold=True, latin=KR)),
             (body, dict(size=9.5, color=DGRAY, latin=KR))],
            line_spacing=1.5, space_after=7)
hairline(s, MX, py, CW, DLINE)
textbox(s, W - MX - Inches(1.2), H - Inches(0.6), Inches(1.2), Inches(0.3),
        [('11', dict(size=9, color=DLABEL, spacing=220))], align=PP_ALIGN.RIGHT)

# ================================================== 12 — Studio (하프 블리드)
s = new_slide()
img_placeholder(s, Inches(7.2), 0, W - Inches(7.2), H, '스튜디오 공간 컷으로 교체\n(풀블리드)', border=False)
textbox(s, MX, Inches(0.34), Inches(6), Inches(0.3),
        [('STUDIO', dict(size=9, color=GRAY, spacing=220))])
hairline(s, MX, Inches(0.72), Inches(6.1))
textbox(s, MX, Inches(1.9), Inches(6.2), Inches(2.4),
        [('SEOUL,', dict(size=58, color=INK, spacing=150)),
         ('YEOKSAM.', dict(size=58, color=INK, spacing=150))],
        line_spacing=1.1)
textbox(s, MX, Inches(4.15), Inches(6.1), Inches(0.6),
        [('서울 강남구 역삼동 783-2, B1', dict(size=16, color=INK, bold=True, latin=KR))])
textbox(s, MX, Inches(4.85), Inches(5.9), Inches(1.8),
        [('제품 · 뷰티 · F&B 촬영에 맞춘 세팅으로 운영하며, 컨셉에 따라 로케이션 촬영도 진행합니다.',
          dict(size=10.5, color=GRAY, latin=KR)),
         ('방문 상담은 이메일 또는 문의 폼으로 일정을 잡아주세요.', dict(size=10.5, color=GRAY, latin=KR))],
        line_spacing=1.6, space_after=8)
textbox(s, MX, H - Inches(0.6), Inches(1.2), Inches(0.3),
        [('12', dict(size=9, color=GRAY, spacing=220))])

# ================================================== 13 — Contact (블랙)
s = new_slide(DARK)
textbox(s, MX, Inches(0.34), Inches(9), Inches(0.3),
        [('CONTACT', dict(size=9, color=DLABEL, spacing=220))])
textbox(s, MX - Inches(0.04), Inches(1.35), Inches(13.5), Inches(4.0),
        [('START', dict(size=120, color=PAPER, spacing=150)),
         ('A PROJECT →', dict(size=120, color=PAPER, spacing=150))],
        line_spacing=1.02)
cy = Inches(5.9)
hairline(s, MX, cy, CW, DLINE)
contact = [
    ('EMAIL', 'studio.holymolly@gmail.com', DISPLAY, 10.5),
    ('INSTAGRAM · THREADS', '@studio_holymolly', DISPLAY, 11.5),
    ('PORTFOLIO', 'studioholymolly.myportfolio.com', DISPLAY, 9.5),
    ('ADDRESS', '역삼동 783-2, B1', KR, 11.5),
]
col_w = CW / 4
for i, (lab, val, lat, vs) in enumerate(contact):
    x = MX + col_w * i
    textbox(s, x, cy + Inches(0.28), col_w - Inches(0.35), Inches(1.0),
            [(lab, dict(size=8.5, color=DLABEL, spacing=200)),
             (val, dict(size=vs, color=PAPER, bold=True, latin=lat))],
            line_spacing=1.8)

# ================================================== 14 — 백커버
s = new_slide()
textbox(s, 0, Inches(3.15), W, Inches(0.7),
        [('STUDIO. HOLYMOLLY', dict(size=22, color=INK, spacing=350))],
        align=PP_ALIGN.CENTER)
textbox(s, 0, Inches(4.05), W, Inches(0.4),
        [('© 2026 — 서울 강남구 역삼동 783-2, B1', dict(size=9, color=GRAY3, latin=KR, spacing=100))],
        align=PP_ALIGN.CENTER)

out = 'STUDIO-HOLYMOLLY_Profile_2026.pptx'
prs.save(out)
print('saved', out, '| slides:', len(prs.slides._sldIdLst))
